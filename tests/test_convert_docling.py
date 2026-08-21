"""The docling engine: PDFs (digital and scanned), decks and workbooks.

Skipped whole unless the `pdf` extra is installed **and** the layout/table
weights are already on disk — the alternative is a test suite that downloads
two thirds of a gigabyte. The weights are looked for where the rest of enough
looks, through the `ENOUGH_WEIGHTS_DIR` seam (`~/enough/weights` by default),
read once at import time before any fixture redirects it. So the house QA
recipe — every `ENOUGH_*` pointed into a scratch dir — exercises this file,
and a machine without the extra skips it instead of failing.

Fixtures are generated, never checked in, for the same reason the `.docx` in
`conftest.py` is: everything needed to build them (Pillow, python-pptx,
openpyxl, and enough's own md → PDF export) is a dependency of the thing
under test, so any machine that can run this file can make them, and a
generated fixture cannot drift away from the reader that has to read it back.

Conversions are session-scoped. Each one pays a fresh worker's torch import
and model load, so converting the same four documents per-test would trade a
minute of wall clock for nothing.
"""

from __future__ import annotations

import os
import re
from pathlib import Path

import pytest

from enough import convert

# Read before any fixture moves it: this is the developer's / harness's real
# weights dir, and the only place a populated artifacts dir can be found
# without downloading one.
_AMBIENT_WEIGHTS = os.environ.get("ENOUGH_WEIGHTS_DIR")


def _models_dir() -> Path | None:
    base = Path(_AMBIENT_WEIGHTS).expanduser() if _AMBIENT_WEIGHTS \
        else Path.home() / "enough" / "weights"
    d = base / "docling"
    try:
        return d if d.is_dir() and any(d.iterdir()) else None
    except OSError:
        return None


MODELS = _models_dir()

pytestmark = pytest.mark.skipif(
    not convert.docling_installed() or MODELS is None,
    reason="needs the pdf extra and prefetched docling weights "
           "(set ENOUGH_WEIGHTS_DIR at a dir whose docling/ is populated)")

OCR_LINES = ["ACME CORPORATION", "Minutes of the annual meeting.",
             "The board approved the budget by a vote of seven to two."]


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture(scope="session", autouse=True)
def docling_env(tmp_path_factory: pytest.TempPathFactory):
    """Scratch for everything convert.py writes globally, real weights dir.

    `ENOUGH_WEIGHTS_DIR` is deliberately *not* redirected: the models are the
    one thing these tests can't fabricate, and nothing here writes into it."""
    state = tmp_path_factory.mktemp("docling-state")
    old = os.environ.get("ENOUGH_EXTRAS_STATE")
    os.environ["ENOUGH_EXTRAS_STATE"] = str(state / "extras.json")
    if MODELS is not None:
        os.environ["ENOUGH_WEIGHTS_DIR"] = str(MODELS.parent)
    convert.reset_engines()
    yield
    if old is None:
        os.environ.pop("ENOUGH_EXTRAS_STATE", None)
    else:
        os.environ["ENOUGH_EXTRAS_STATE"] = old
    convert.reset_engines()


def _write_pdf(md_text: str, dest: Path, images: dict[str, bytes] | None = None) -> Path:
    """markdown → PDF through enough's own export path, so the fixture is made
    by the code the reader is paired with."""
    from enough import convert_worker
    work = dest.parent
    for name, blob in (images or {}).items():
        (work / name).write_bytes(blob)
    src = work / f"_src-{dest.stem}.md"
    src.write_text(md_text, encoding="utf-8")
    convert_worker.do_export({"twin": str(src), "out": str(dest), "target": ".pdf",
                              "resource_path": str(work), "reference_doc": None})
    src.unlink()
    for name in (images or {}):
        (work / name).unlink()
    return dest


def _plate_png() -> bytes:
    """A photographic-looking plate. A line drawing or a box of text gets
    classified as text by the layout model and never becomes a picture."""
    import io
    import math

    from PIL import Image
    im = Image.new("RGB", (320, 200))
    px = im.load()
    for y in range(200):
        for x in range(320):
            px[x, y] = (int(127 + 120 * math.sin(x / 17.0)),
                        int(127 + 120 * math.sin(y / 11.0)),
                        int(127 + 120 * math.sin((x + y) / 23.0)))
    buf = io.BytesIO()
    im.save(buf, "PNG")
    return buf.getvalue()


@pytest.fixture(scope="session")
def docs(tmp_path_factory: pytest.TempPathFactory) -> dict[str, Path]:
    """One project dir with a digital PDF, a scanned PDF, a deck and a
    workbook, each already converted."""
    proj = tmp_path_factory.mktemp("docling-proj")

    _write_pdf(
        "# Quarterly Report\n\n"
        "![the plate](plate.png)\n\n"
        "## The numbers\n\n"
        "| Region | Units | Margin |\n|---|---|---|\n"
        "| North | 512 | 31% |\n| South | 402 | 28% |\n\n"
        "A closing paragraph, so the table is not the last thing on the page.\n",
        proj / "report.pdf", {"plate.png": _plate_png()})

    # Image-only: no text layer at all, so anything that comes out is OCR.
    from PIL import Image, ImageDraw, ImageFont
    page = Image.new("RGB", (1275, 1650), "white")
    draw = ImageDraw.Draw(page)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial.ttf", 44)
    except OSError:
        font = ImageFont.load_default(size=44)
    y = 220
    for line in OCR_LINES:
        draw.text((150, y), line, fill="black", font=font)
        y += 100
    page.save(proj / "scan.pdf", "PDF", resolution=150.0)

    from pptx import Presentation
    prs = Presentation()
    first = prs.slides.add_slide(prs.slide_layouts[0])
    first.shapes.title.text = "Widget Division Review"
    first.placeholders[1].text = "Second quarter"
    second = prs.slides.add_slide(prs.slide_layouts[1])
    second.shapes.title.text = "Highlights"
    frame = second.placeholders[1].text_frame
    frame.text = "Units up twelve percent"
    frame.add_paragraph().text = "Backlog cleared"
    prs.save(proj / "deck.pptx")

    from openpyxl import Workbook
    book = Workbook()
    sheet = book.active
    sheet.title = "Regions"
    sheet.append(["Region", "Units"])
    sheet.append(["North", 512])
    sheet.append(["South", 402])
    book.save(proj / "book.xlsx")

    out = {p.stem: p for p in (proj / "report.pdf", proj / "scan.pdf",
                               proj / "deck.pptx", proj / "book.xlsx")}
    for original in out.values():
        convert.run_convert(original)
    return out


def twin_text(original: Path) -> str:
    return convert.twin_path(original).read_text(encoding="utf-8")


# ---------------------------------------------------------------------------
# The engine reports itself
# ---------------------------------------------------------------------------

def test_the_engine_is_available_once_the_models_are_there(docling_env):
    assert convert.docling_models_present() is True
    assert convert.docling_available() is True
    assert convert.engine_available("docling") is True
    row = next(r for r in convert.formats_view()["formats"] if r["ext"] == ".pdf")
    assert row["available"] is True


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def test_a_digital_pdf_becomes_a_structured_twin(docs):
    text = twin_text(docs["report"])
    assert re.search(r"^#+ .*Quarterly Report", text, re.M)
    assert "The numbers" in text
    # The table survives as a markdown table, not as prose or raw HTML.
    rows = [ln for ln in text.splitlines() if ln.strip().startswith("|")]
    assert len(rows) >= 4
    assert "North" in text and "512" in text
    assert "<table" not in text.lower()


def test_the_twin_is_plain_markdown_with_no_absolute_paths(docs):
    for name in ("report", "scan", "deck", "book"):
        text = twin_text(docs[name])
        assert "<img" not in text.lower(), name
        assert "<!--" not in text, name          # docling's image placeholder
        assert "/private/" not in text, name     # save_as_markdown writes abs paths
        assert not text.startswith("---"), name  # no front matter (plan §2)


def test_a_scanned_pdf_is_read_by_ocr(docs):
    text = twin_text(docs["scan"])
    # The fixture is an image with no text layer, so every word below can only
    # have come out of the OCR stage.
    assert "ACME CORPORATION" in text
    assert "annual meeting" in text
    manifest = convert.read_manifest(docs["scan"])
    assert manifest["engine"]["name"] == "docling"
    assert manifest["engine"]["ocr"] in ("ocrmac", "rapidocr")
    assert manifest["engine"]["version"]


def test_a_pdfs_pictures_land_in_the_pandoc_asset_layout(docs):
    original = docs["report"]
    assets = convert.assets_dir(original)
    assert assets.is_dir(), "the plate should have been extracted"
    files = sorted(p.name for p in assets.iterdir())
    assert files == ["img-1.png"]
    text = twin_text(original)
    # Flat, relative to the twin, and the exact shape the pandoc reader emits
    # so renderMarkdown's blob rewrite needs no second case.
    assert f"]({assets.name}/img-1.png)" in text
    assert convert.read_manifest(original)["assets"] == assets.name


def test_a_document_with_no_pictures_gets_no_assets_dir(docs):
    assert not convert.assets_dir(docs["scan"]).exists()
    assert convert.read_manifest(docs["scan"])["assets"] is None


# ---------------------------------------------------------------------------
# Decks and workbooks
# ---------------------------------------------------------------------------

def test_a_deck_becomes_headed_sections(docs):
    text = twin_text(docs["deck"])
    assert re.search(r"^#+ .*Widget Division Review", text, re.M)
    assert re.search(r"^#+ .*Highlights", text, re.M)
    assert "Backlog cleared" in text
    # No OCR stage runs for a deck, and the manifest says so rather than
    # naming an engine that never ran.
    assert convert.read_manifest(docs["deck"])["engine"]["ocr"] is None


def test_a_workbook_becomes_markdown_tables(docs):
    text = twin_text(docs["book"])
    assert "| Region" in text and "| North" in text
    assert "512" in text and "402" in text
    assert convert.read_manifest(docs["book"])["engine"]["name"] == "docling"


# ---------------------------------------------------------------------------
# Lifecycle
# ---------------------------------------------------------------------------

def test_a_converted_pdf_reads_fresh_then_edited(docs):
    original = docs["report"]
    assert convert.state(original) == "fresh"
    twin = convert.twin_path(original)
    twin.write_text(twin.read_text(encoding="utf-8") + "\nan edit.\n",
                    encoding="utf-8")
    assert convert.state(original) == "edited"


def test_reconvert_keeps_docling_asset_names_stable(tmp_path: Path):
    """The assets dir is cleared before each conversion, so the same PDF
    always lands on `img-1.png` instead of accumulating collision suffixes —
    the same guarantee the pandoc reader gives."""
    pdf = _write_pdf("# Plate\n\n![p](plate.png)\n\nAfter the plate.\n",
                     tmp_path / "plate-doc.pdf", {"plate.png": _plate_png()})
    convert.run_convert(pdf)
    assets = convert.assets_dir(pdf)
    assert sorted(p.name for p in assets.iterdir()) == ["img-1.png"]
    first = (assets / "img-1.png").read_bytes()

    convert.run_convert(pdf, force=True)
    assert sorted(p.name for p in assets.iterdir()) == ["img-1.png"]
    assert (assets / "img-1.png").read_bytes() == first
    assert f"]({assets.name}/img-1.png)" in twin_text(pdf)


def test_the_worker_refuses_a_format_docling_does_not_read(tmp_path: Path):
    from enough import convert_worker
    odt = tmp_path / "notes.odt"
    odt.write_bytes(b"x")
    with pytest.raises(convert_worker.WorkerError) as e:
        convert_worker.do_convert({
            "op": "convert", "engine": "docling", "original": str(odt),
            "twin": str(tmp_path / "notes.odt.md"),
            "assets": str(tmp_path / "notes.odt.assets"),
            "artifacts_dir": str(MODELS)})
    assert ".odt" in str(e.value)
